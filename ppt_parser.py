import os
from pptx import Presentation
from base_parser import BaseParser
from typing import List, Dict, Any


class PPTParser(BaseParser):
    """Parser for PPT/PPTX files using python-pptx library - text only"""
    
    def __init__(self):
        super().__init__("ppt")
    
    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract text from PPTX file using your working code logic"""
        try:
            presentation = Presentation(pptx_path)
            extracted_text = ""
            
            for slide_number, slide in enumerate(presentation.slides):
                extracted_text += f"\nSlide {slide_number + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
            
            return extracted_text.strip()
        
        except Exception as e:
            print(f"Error extracting text from {pptx_path}: {e}")
            return ""
    
    def extract_sections(self, file_path: str) -> List[str]:
        """Extract sections from PPT/PPTX file (each slide as a section)"""
        try:
            # Handle .ppt extension by treating it as .pptx
            original_path = file_path
            if file_path.lower().endswith('.ppt'):
                # Change extension to .pptx for processing
                file_path = file_path[:-4] + '.pptx'
                # If the .pptx version doesn't exist, try with original .ppt
                if not os.path.exists(file_path):
                    file_path = original_path
            
            presentation = Presentation(file_path)
            sections = []
            
            # Process each slide as a separate section
            for slide_number, slide in enumerate(presentation.slides, 1):
                slide_content = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if len(text) > 0:  # Keep all text content
                            slide_content.append(text)
                
                # Combine all content from this slide
                if slide_content:
                    section_text = f"Slide {slide_number}: " + " ".join(slide_content)
                    sections.append(section_text)
            
            # If no sections found, fall back to full text extraction
            if not sections:
                full_text = self.extract_text_from_pptx(file_path)
                if full_text:
                    # Split by slides for sectioning
                    slide_texts = full_text.split('\nSlide ')
                    for i, slide_text in enumerate(slide_texts[1:], 1):  # Skip first empty split
                        if slide_text.strip():
                            sections.append(f"Slide {i}: {slide_text.strip()}")
            
            return sections
        
        except Exception as e:
            print(f"Error extracting sections from {file_path}: {e}")
            return ["No text could be extracted from presentation."]
    
    def parse_file(self, file_path: str, uploaded_by: str = "") -> Dict[str, Any]:
        """Parse a PPT/PPTX file and return structured + unstructured format"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            # Extract sections (slide text content)
            sections = self.extract_sections(file_path)
            
            # Generate file ID once for consistency
            file_id = self.generate_file_id()
            
            # Create structured data (each slide as separate post)
            structured = []
            for section in sections:
                structured.append(self.create_post(file_id, section))
            
            # Create unstructured data (full combined text)
            full_text = " ".join(sections)
            unstructured = [self.create_post(file_id, full_text)] if full_text else []
            
            # Return in the requested format
            return {
                "structured": structured,
                "unstructured": unstructured
            }
            
        except Exception as e:
            return {
                "structured": [],
                "unstructured": [{
                    'error': str(e),
                    'file_path': file_path,
                    'source_type': self.source_type
                }]
            }
